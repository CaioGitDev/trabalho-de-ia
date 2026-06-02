"""Gera o PowerPoint completo para o video-relatorio.

Cria docs/apresentacao.pptx com um slide por seccao do guiao
(docs/guiao_video_para_ler.md), incorporando o diagrama do pipeline e o
grafico de comparacao de modelos.

Executar:
    pip install python-pptx
    python src/gerar_powerpoint.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "docs" / "figures"
OUT = ROOT / "docs" / "apresentacao.pptx"

# paleta
AZUL = RGBColor(0x1F, 0x4E, 0x79)
AZUL_CLARO = RGBColor(0x3B, 0x7D, 0xD8)
VERDE = RGBColor(0x2C, 0xA0, 0x2C)
VERMELHO = RGBColor(0xD6, 0x27, 0x28)
LARANJA = RGBColor(0xF0, 0x8C, 0x00)
CINZA = RGBColor(0x55, 0x55, 0x55)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
ESCURO = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_band(slide, color=AZUL, height=Inches(1.1)):
    """Faixa de cor no topo."""
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def set_run(run, text, size, color=ESCURO, bold=False, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"


def title_in_band(slide, text, color=AZUL):
    add_band(slide, color)
    tf = textbox(slide, Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, 30, BRANCO, bold=True)


def bullets(slide, items, top=Inches(1.5), left=Inches(0.7),
            width=Inches(12.0), height=Inches(5.5), size=20):
    """items: lista de (texto, nivel, cor, bold)."""
    tf = textbox(slide, left, top, width, height)
    first = True
    for it in items:
        text, level, color, bold = (it + (0, ESCURO, False))[:4] if isinstance(it, tuple) else (it, 0, ESCURO, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(10)
        prefix = "" if level == 0 else "– "
        set_run(p.add_run(), prefix + text, size, color, bold=bold)


def add_image_centered(slide, path, top, max_w, max_h):
    from PIL import Image
    with Image.open(path) as im:
        w, h = im.size
    ratio = min(max_w / w, max_h / h)
    nw, nh = int(w * ratio), int(h * ratio)
    left = int((SW - nw) / 2)
    slide.shapes.add_picture(str(path), left, top, width=nw, height=nh)


# ---------------------------------------------------------------- Slide 1: titulo
s = add_slide()
from pptx.enum.shapes import MSO_SHAPE
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = AZUL; bg.line.fill.background(); bg.shadow.inherit = False

tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Análise de Rótulos Alimentares", 44, BRANCO, bold=True)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
set_run(p2.add_run(), "Deteção de glúten e alergénios com IA explicável", 24, RGBColor(0xCF, 0xE0, 0xF5))

tf2 = textbox(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.5))
p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "[Nome 1]  ·  [Nome 2]  ·  [Nome 3]", 20, BRANCO, bold=True)
p2 = tf2.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
set_run(p2.add_run(), "[Disciplina]  —  Universidade Fernando Pessoa", 16, RGBColor(0xCF, 0xE0, 0xF5))

# ---------------------------------------------------------------- Slide 2: problema
s = add_slide()
title_in_band(s, "O problema")
bullets(s, [
    ("Para celíacos e alérgicos, ler um rótulo é difícil e demorado.", 0, ESCURO, True),
    ("Nomes técnicos de ingredientes (amido modificado, maltodextrina…).", 0, CINZA, False),
    ("Letra pequena e rótulos em várias línguas.", 0, CINZA, False),
    ("Um erro de leitura pode ter consequências de saúde.", 0, CINZA, False),
    ("A nossa proposta: um sistema de IA que faz esta leitura a partir de uma fotografia do rótulo.", 0, AZUL, True),
], size=22)

# ---------------------------------------------------------------- Slide 3: objetivo
s = add_slide()
title_in_band(s, "Objetivo do sistema")
bullets(s, [
    ("A partir de uma foto (ou texto) dos ingredientes, devolve:", 0, ESCURO, True),
    ("1. Deteção da presença de glúten.", 1, ESCURO, False),
    ("2. Identificação dos 14 alergénios da UE (Reg. 1169/2011).", 1, ESCURO, False),
    ("3. Sinalização de ingredientes ambíguos.", 1, ESCURO, False),
    ("4. Grau de risco para o consumidor.", 1, ESCURO, False),
    ("Ponto-chave: resultado EXPLICÁVEL — justificação ingrediente a ingrediente.", 0, VERDE, True),
], size=22)

# ---------------------------------------------------------------- Slide 4: arquitetura (diagrama)
s = add_slide()
title_in_band(s, "Arquitetura — pipeline")
pipe = FIG / "00_pipeline.png"
if pipe.exists():
    add_image_centered(s, pipe, Inches(1.4), Inches(12.0), Inches(3.4))
tf = textbox(s, Inches(0.7), Inches(5.4), Inches(12.0), Inches(1.6))
p = tf.paragraphs[0]
set_run(p.add_run(), "OCR → NLP → Regras + ML → Explicação → Relatório → Avaliação", 20, AZUL, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "Reutiliza conceitos das aulas: Levenshtein/Jaccard, Naive Bayes, Árvores, Random Forest, matriz de confusão.", 16, CINZA)

# ---------------------------------------------------------------- Slide 5: porque nao LLM
s = add_slide()
title_in_band(s, "Porquê não um LLM?")
bullets(s, [
    ("Objetivo do trabalho: aplicar os conceitos dados nas aulas.", 0, ESCURO, True),
    ("Explicabilidade: um LLM seria uma caixa preta — não justificaria cada decisão.", 0, ESCURO, True),
    ("Custo e simplicidade: corre em CPU, sem dependências pesadas.", 0, ESCURO, True),
    ("A justificação por ingrediente é o requisito central — e exige regras transparentes.", 0, AZUL, True),
], size=22)

# ---------------------------------------------------------------- Slide 6: demo (placeholder)
s = add_slide()
title_in_band(s, "Demonstração — aplicação web", AZUL_CLARO)
bullets(s, [
    ("Aplicação Flask: cola o texto dos ingredientes ou carrega a foto do rótulo.", 0, ESCURO, True),
    ("Caso com glúten → veredicto \"Contém Glúten\" (vermelho); ❌ no trigo/cevada/aveia.", 0, VERMELHO, False),
    ("Caso seguro → veredicto \"Sem Glúten\" (verde).", 0, VERDE, False),
    ("Caso ambíguo → \"Médio Risco\" (amarelo): sinaliza a dúvida em vez de assumir.", 0, LARANJA, False),
    ("Cada ingrediente tem a sua justificação — análise explicável de ponta a ponta.", 0, AZUL, True),
], size=21)
tf = textbox(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.7))
set_run(tf.paragraphs[0].add_run(),
        "(No vídeo: gravação de ecrã da app a correr — http://127.0.0.1:5000)", 14, CINZA, italic=True)

# ---------------------------------------------------------------- Slide 7: resultados
s = add_slide()
title_in_band(s, "Resultados — avaliação")
bullets(s, [
    ("A) Modelos de ML comparados — melhor: Random Forest", 0, ESCURO, True),
    ("Exatidão 88,7% (teste). RF > Árvore de Decisão > Naive Bayes.", 1, CINZA, False),
    ("Limitação: rótulos gerados pelas próprias regras (supervisão fraca).", 1, CINZA, False),
    ("B) Avaliação independente e NÃO circular", 0, ESCURO, True),
    ("Pipeline lê só os ingredientes; compara com o campo 'allergens' da OFF.", 1, CINZA, False),
    ("Exatidão 91,2%  ·  recall de glúten ~90% (subiu de 80% com léxico EN/FR).", 1, VERDE, True),
], size=20)

# ---------------------------------------------------------------- Slide 8: grafico modelos
s = add_slide()
title_in_band(s, "Comparação de modelos")
graf = FIG / "08_model_comparison.png"
if graf.exists():
    add_image_centered(s, graf, Inches(1.35), Inches(11.5), Inches(5.4))

# ---------------------------------------------------------------- Slide 9: limitacoes
s = add_slide()
title_in_band(s, "Limitações (assumidas)")
bullets(s, [
    ("OCR é a maior fonte de erro: ~48% de confiança em fotos reais.", 0, ESCURO, True),
    ("Funciona melhor sobre texto de ingredientes limpo.", 1, CINZA, False),
    ("Ground-truth ruidoso: a OFF omite alergénios — a precisão real é por vezes superior.", 0, ESCURO, True),
    ("Opção deliberada: assinalar risco a mais (mais seguro para o consumidor).", 0, ESCURO, True),
    ("Línguas cobertas: PT + EN + FR (extensão natural: IT/ES).", 0, CINZA, False),
], size=21)

# ---------------------------------------------------------------- Slide 10: conclusao
s = add_slide()
title_in_band(s, "Conclusão", VERDE)
bullets(s, [
    ("Cumprimos todos os objetivos do enunciado.", 0, ESCURO, True),
    ("Sistema completo: da imagem ao relatório final.", 0, CINZA, False),
    ("Deteta glúten e os 14 alergénios, atribui grau de risco.", 0, CINZA, False),
    ("Acima de tudo: justifica cada decisão de forma explicável.", 0, VERDE, True),
    ("Aplicação web funcional como concretização do enunciado.", 0, CINZA, False),
], size=22)

# ---------------------------------------------------------------- Slide 11: fim
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = AZUL; bg.line.fill.background(); bg.shadow.inherit = False
tf = textbox(s, Inches(1.0), Inches(2.8), Inches(11.3), Inches(2.0))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Obrigado pela atenção", 40, BRANCO, bold=True)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
set_run(p2.add_run(), "[Nome 1]  ·  [Nome 2]  ·  [Nome 3]", 22, RGBColor(0xCF, 0xE0, 0xF5))

# ---------------------------------------------------------------- Notas de orador
# Texto do guiao (docs/guiao_video_para_ler.md), uma entrada por slide, pela ordem.
NOTAS = [
    # 1 - titulo
    "Olá. Este é o nosso trabalho de [disciplina], desenvolvido por [Nome 1], "
    "[Nome 2] e [Nome 3].",
    # 2 - problema
    "O nosso projeto responde a um problema do dia a dia: para uma pessoa celíaca, "
    "ou com alergias alimentares, ler o rótulo de um produto e perceber se o pode "
    "comer é difícil e demorado. Os ingredientes têm nomes técnicos, estão em letras "
    "pequenas, e muitas vezes em várias línguas. A nossa proposta é um sistema de "
    "inteligência artificial que faz esse trabalho automaticamente, a partir de uma "
    "simples fotografia do rótulo.",
    # 3 - objetivo
    "O sistema recebe a foto de um rótulo, ou o texto dos ingredientes, e devolve "
    "quatro coisas. Primeiro, deteta a presença de glúten. Segundo, identifica quais "
    "dos catorze alergénios obrigatórios da União Europeia estão presentes — os do "
    "Regulamento mil cento e sessenta e nove de dois mil e onze. Terceiro, assinala "
    "os ingredientes ambíguos, que podem ou não conter glúten consoante a origem. E "
    "quarto, atribui um grau de risco. Mas o ponto mais importante é que o resultado "
    "é explicável: o sistema não dá apenas um veredicto, justifica-o ingrediente a "
    "ingrediente.",
    # 4 - arquitetura
    "Construímos um pipeline com seis etapas, que reutiliza os conceitos das aulas. "
    "A primeira é o OCR, que converte a imagem em texto com o Tesseract e "
    "pré-processamento em OpenCV. A segunda é o processamento de linguagem natural, "
    "que limpa o texto, separa os ingredientes e faz correspondência com o léxico "
    "usando distâncias como Levenshtein e Jaccard. A terceira é a classificação, que "
    "combina regras com modelos de aprendizagem automática. A quarta é o motor de "
    "explicação. A quinta gera o relatório final. E a sexta é a avaliação.",
    # 5 - porque nao LLM
    "Uma pergunta natural é: porque não usámos um modelo de linguagem grande, como o "
    "ChatGPT? Por duas razões. Primeiro, porque o objetivo do trabalho era aplicar os "
    "conceitos das aulas. E segundo, porque um modelo desses seria uma caixa preta — "
    "não conseguiríamos justificar cada decisão, que é precisamente o requisito "
    "central do projeto.",
    # 6 - demo
    "Vamos ver o sistema a funcionar na aplicação web. Começo por colar a lista de "
    "ingredientes de um produto com glúten: aveia, malte de cevada, farinha de trigo, "
    "açúcar, óleo e sal. O resultado é imediato: o veredicto 'Contém Glúten', a "
    "vermelho. Na tabela, a aveia, a cevada e o trigo aparecem com cruz vermelha, e "
    "cada uma diz porquê. O açúcar, o óleo e o sal estão a verde. Agora um produto "
    "seguro: água, tomate, cenoura, cebola, azeite, sal. O veredicto muda para 'Sem "
    "Glúten', a verde — é o mesmo motor, não está a adivinhar. E quando os "
    "ingredientes são ambíguos, atribui risco médio e sinaliza a dúvida.",
    # 7 - resultados
    "Avaliámos o sistema de duas formas. Na primeira, comparámos os modelos de "
    "aprendizagem automática — o melhor foi o Random Forest, com oitenta e oito "
    "vírgula sete por cento de exatidão. Mas como os rótulos de treino foram gerados "
    "pelas próprias regras, esta avaliação mede sobretudo a capacidade de reproduzir "
    "essa lógica. Por isso fizemos uma segunda avaliação, independente e não circular: "
    "o pipeline lê só o texto dos ingredientes e comparamos com o campo de alergénios "
    "da Open Food Facts, que o sistema nunca consulta. Aqui obtivemos noventa e um "
    "vírgula dois por cento de exatidão, com deteção de glúten de quase noventa por "
    "cento — que subiu de oitenta quando alargámos o léxico ao inglês e francês.",
    # 8 - grafico modelos
    "Este gráfico mostra a comparação entre os três modelos nas várias métricas. O "
    "Random Forest fica à frente em exatidão e F1, seguido da Árvore de Decisão e do "
    "Naive Bayes. A classe mais difícil para todos é a dos ingredientes 'suspeitos', "
    "ou ambíguos.",
    # 9 - limitacoes
    "Quanto às limitações, somos transparentes. A maior fonte de erro é o OCR: em "
    "fotografias reais a confiança ronda os quarenta e oito por cento e parte do texto "
    "perde-se, por isso o sistema funciona melhor sobre texto limpo. A base de dados "
    "também tem ruído — alguns produtos com glúten não o declaram — o que significa "
    "que a precisão real é por vezes superior à medida. E, por opção, preferimos "
    "assinalar risco a mais do que a menos: é a escolha mais segura.",
    # 10 - conclusao
    "Em conclusão: cumprimos todos os objetivos do enunciado. Construímos um sistema "
    "completo, da imagem ao relatório, que deteta glúten e alergénios, atribui um grau "
    "de risco e, acima de tudo, justifica cada decisão de forma explicável.",
    # 11 - fim
    "Obrigado pela atenção.",
]

for slide, nota in zip(prs.slides, NOTAS):
    slide.notes_slide.notes_text_frame.text = nota

prs.save(str(OUT))
print(f"Gravado: {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
      f"{len(NOTAS)} notas de orador)")
